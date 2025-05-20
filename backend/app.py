import os
from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from dotenv import load_dotenv
import requests
import speech_recognition as sr
from gtts import gTTS
import pygame
import time
from groq import Groq
import random
import string
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from datetime import datetime, timedelta
import logging
from logging.handlers import RotatingFileHandler
import pymongo
from bson.objectid import ObjectId
import ssl
import io
import csv
from io import BytesIO
import json
import re
from PIL import Image
import pytesseract

# Configure logging first - before any logger references
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)

# Create console handler only (removing file handler)
stream_handler = logging.StreamHandler()

# Create formatter and add it to handler
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
stream_handler.setFormatter(formatter)

# Add handler to the logger
logger.addHandler(stream_handler)

# Now check for optional imports
try:
    import pandas as pd
    pandas_available = True
except ImportError:
    pandas_available = False
    logger.warning("pandas not installed - Excel/CSV processing will be limited")

try:
    import docx
    docx_available = True
except ImportError:
    docx_available = False
    logger.warning("python-docx not installed - Word document processing will be limited")

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False
    logger.warning("python-pptx not installed - PowerPoint processing will be limited")

try:
    import pytesseract
    tesseract_available = True
except ImportError:
    tesseract_available = False
    logger.warning("pytesseract not installed - Image OCR will be limited")

#load_dotenv('/etc/bharatai.env')
load_dotenv('.env')

# Function to create directory with proper permissions for VPS
def create_directory_with_permissions(dir_path):
    """Create directory if it doesn't exist and set proper permissions for VPS"""
    if not os.path.exists(dir_path):
        os.makedirs(dir_path, exist_ok=True)
        try:
            # Set directory permissions to 755 for VPS environments
            os.chmod(dir_path, 0o755)
            logger.info(f"Created directory with proper permissions: {dir_path}")
        except Exception as e:
            logger.warning(f"Could not set permissions for {dir_path}: {str(e)}")

# Create required directories with proper permissions (consolidate all directory creations here)
create_directory_with_permissions("static")
create_directory_with_permissions("temp")
create_directory_with_permissions("pdf_files")
create_directory_with_permissions("audio_files")
create_directory_with_permissions("training_data")
create_directory_with_permissions("trained_models")
create_directory_with_permissions("embedding_models")
create_directory_with_permissions("cached_embeddings")

def get_database():
    """
    Initialize and return MongoDB connection
    """
    # Get MongoDB URI from environment variables
    mongo_url = os.getenv('MONGO_URL') or os.getenv('MONGO_URI')
    
    if not mongo_url:
        logger.error("MongoDB URI not found in environment variables")
        return None
    
    try:
        # Set MongoDB connection options
        # Updated to use newer PyMongo SSL options
        mongo_options = {
            'tls': True,
            'tlsAllowInvalidCertificates': True,
            'retryWrites': False,
            'connectTimeoutMS': 5000,
            'socketTimeoutMS': 30000
        }
        
        # Create a MongoDB client with the options
        client = pymongo.MongoClient(mongo_url, **mongo_options)
        # Test the connection
        client.admin.command('ping')
        logger.info("Successfully connected to MongoDB")
        
        # Initialize collections if they don't exist
        db = client.get_database("chatbotDB")
        if "users" not in db.list_collection_names():
            db.create_collection("users")
            db.users.create_index([("email", pymongo.ASCENDING)], unique=True)
            logger.info("Created users collection")
            
        if "verifications" not in db.list_collection_names():
            db.create_collection("verifications")
            db.verifications.create_index([("email", pymongo.ASCENDING)])
            db.verifications.create_index([("code", pymongo.ASCENDING)])
            logger.info("Created verifications collection")
            
        if "chats" not in db.list_collection_names():
            db.create_collection("chats")
            db.chats.create_index([("user_id", pymongo.ASCENDING)])
            db.chats.create_index([("updated_at", pymongo.DESCENDING)])
            logger.info("Created chats collection")
            
        return db
    except Exception as e:
        logger.error(f"Failed to connect to MongoDB: {str(e)}")
        return None
    
# Get database connection
db = get_database()
if db is None:  # Fixed comparison to check if db is None
    logger.critical("Could not connect to MongoDB. Application will not function correctly.")

# Initialize the Flask app
app = Flask(__name__)

# Configure CORS
CORS(app, resources={r"/api/*": {"origins": "*", "supports_credentials": True}})

# Initialize Groq client
groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))

# Supported languages
LANGUAGES = {
    "english": "en",
    "hindi": "hi",
    "kannada": "kn",
    "tamil": "ta",
    "telugu": "te",
    "sanskrit": "sa"
}

# Initialize speech recognition
recognizer = sr.Recognizer()
# Try to initialize speech recognition components
try:
    microphone = sr.Microphone()
    speech_recognition_available = True
except (ImportError, AttributeError):
    logger.warning("PyAudio not found. Speech recognition features will be disabled.")
    microphone = None
    speech_recognition_available = False

# Model mapping
MODEL_MAPPING = {
    "LLaMA3-versatile": "llama-3.3-70b-versatile",
    "LLaMA3": "llama3-70b-8192",
    "LLaMA2": "llama2-70b-4096"
}

# Helper Functions
def generate_verification_code():
    """Generate a 6-digit verification code"""
    return ''.join(random.choices(string.digits, k=6))

def send_verification_email(email, code):
    """Send verification email with OTP code"""
    sender_email = os.getenv("EMAIL_SENDER", "noreply@chatbotapp.com")
    sender_password = os.getenv("EMAIL_PASSWORD", "")

    logger.info(f"Sending verification code {code} to {email}")

    try:
        message = MIMEMultipart()
        message["From"] = sender_email
        message["To"] = email
        message["Subject"] = "Email Verification - Bharat AI"

        body = f"""
        <html>
        <body>
            <h2>Verify Your Email Address</h2>
            <p>Thank you for registering with Bharat AI. Use the verification code below to complete your registration:</p>
            <h1 style="color: #3b82f6;">{code}</h1>
            <p>This code will expire in 10 minutes.</p>
            <p>If you didn't request this verification, please ignore this email.</p>
        </body>
        </html>
        """

        message.attach(MIMEText(body, "html"))

        if sender_password:
            with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
                server.login(sender_email, sender_password)
                server.sendmail(sender_email, email, message.as_string())
            logger.info(f"Email sent successfully to {email}")
            return True
        else:
            # Simulate successful sending for development
            logger.warning("Email service not configured. Would have sent:")
            logger.warning(f"To: {email}")
            logger.warning(f"Code: {code}")
            return True
    except Exception as e:
        logger.error(f"Email sending error: {str(e)}")
        return False

def get_groq_response(prompt, language="en", model="llama3-70b-8192"):
    """Get response from Groq API"""
    try:
        completion = groq_client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            max_tokens=1024,
            top_p=1,
            frequency_penalty=0,
            presence_penalty=0,
            stop=None,
        )
        return completion.choices[0].message.content
    except Exception as e:
        logger.error(f"Error with Groq API: {str(e)}")
        return get_error_message(language)

def get_error_message(language):
    """Return error message in appropriate language"""
    messages = {
        "en": "Sorry, I encountered an error. Please try again.",
        "hi": "क्षमा करें, मुझे एक समस्या हुई। कृपया पुनः प्रयास करें।",
        "kn": "ಕ್ಷಮಿಸಿ, ನಾನು ಸಮಸ್ಯೆಯನ್ನು ಎದುರಿಸಿದ್ದೇನೆ. ದಯವಿಟ್ಟು ಮತ್ತೆ ಪ್ರಯತ್ನಿಸಿ.",
        "ta": "மன்னிக்கவும், நான் ஒரு சிக்கலை எதிர்கொண்டேன். தயவு செய்து மீண்டும் முயற்சிக்கவும்.",
        "te": "క్షమించండి, నేను ఒక సమస్యను ఎదుర్కొన్నాను. దయచేసి మళ్లీ ప్రయత్నించండి.",
        "sa": "क्षम्यताम्, अहं समस्याम् अनुभवम्। कृपया पुनः प्रयत्नं कुरुत।"
    }
    return messages.get(language, messages["en"])

# Languages supported for speech recognition
SPEECH_LANGUAGES = ["en", "hi", "kn", "ta", "te"]

def process_audio(audio_data, language="en"):
    """Process audio data and return text"""
    if language not in SPEECH_LANGUAGES:
        logger.warning(f"Language {language} not supported for speech recognition")
        return None

    try:
        text = recognizer.recognize_google(audio_data, language=language)
        logger.info(f"Recognized speech: {text}")
        return text
    except Exception as e:
        logger.error(f"Error in speech recognition: {str(e)}")
        return None

# Languages supported for text-to-speech
TTS_LANGUAGES = ["en", "hi", "kn", "ta", "te"]

def text_to_speech(text, lang="en"):
    """Convert text to speech and return audio file path"""
    if lang not in TTS_LANGUAGES:
        logger.warning(f"Language {lang} not supported for text-to-speech")
        return None

    try:
        tts = gTTS(text=text, lang=lang)
        create_directory_with_permissions("audio_files")
        filename = os.path.join("audio_files", f"response_{random.randint(1000, 9999)}.mp3")
        tts.save(filename)
        return filename
    except Exception as e:
        logger.error(f"Error in speech synthesis: {str(e)}")
        return None

from langdetect import detect

# Define PDF extraction functions at the module level (not inside a try/except block)
def extract_text_with_pdfminer(pdf_path):
    """Extract text from PDF using pdfminer"""
    try:
        from pdfminer.high_level import extract_text
        return extract_text(pdf_path)
    except ImportError:
        logger.warning("pdfminer.six not installed")
        return None

def extract_text_with_pdftotext(pdf_path):
    """Extract text from PDF using pdftotext command line tool"""
    try:
        import subprocess
        result = subprocess.run(["pdftotext", pdf_path, "-"], capture_output=True, text=True)
        if result.returncode == 0:
            return result.stdout
        else:
            logger.warning(f"pdftotext command failed: {result.stderr}")
            return None
    except (subprocess.SubprocessError, FileNotFoundError):
        logger.warning("pdftotext command not available")
        return None

def extract_text_with_pypdf2(pdf_path):
    """Extract text from PDF using PyPDF2"""
    try:
        import PyPDF2
        text = ""
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
        return text
    except ImportError:
        logger.warning("PyPDF2 not installed")
        return None
    except Exception as e:
        logger.warning(f"PyPDF2 extraction error: {str(e)}")
        return None

# Attempt to correctly import and configure PyMuPDF
pdf_support = True  # Start with assumption of PDF support

# PyMuPDF is optional - gracefully handle if missing
try:
    import fitz  # This is PyMuPDF
    def extract_text_with_pymupdf(pdf_path):
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text
    logger.info("Successfully imported PyMuPDF as fitz module")
except ImportError:
    try:
        import PyMuPDF
        fitz = PyMuPDF
        def extract_text_with_pymupdf(pdf_path):
            doc = fitz.open(pdf_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()
            return text
        logger.info("Successfully imported PyMuPDF directly")
    except (ImportError, AttributeError):
        logger.warning("PyMuPDF/fitz not available - will use alternative PDF libraries")
        extract_text_with_pymupdf = lambda path: None

# Set up the main extraction function that will try all available methods
def extract_text_from_pdf(pdf_path):
    """Extract text from PDF using all available methods"""
    # Try each method in order (reordering to prioritize more reliable methods)
    extractors = [
        ("PyPDF2", extract_text_with_pypdf2),  # Move PyPDF2 first since it has fewer dependencies
        ("pdfminer", extract_text_with_pdfminer),  # Then pdfminer which should work reliably
        ("PyMuPDF", extract_text_with_pymupdf),  # Keep PyMuPDF as an option if available
        ("pdftotext", extract_text_with_pdftotext),  # Last option
    ]
    
    for name, extractor in extractors:
        try:
            logger.info(f"Trying PDF extraction with {name}")
            text = extractor(pdf_path)
            if text and text.strip():
                logger.info(f"Successfully extracted text with {name}")
                return text
        except Exception as e:
            logger.warning(f"{name} extraction failed: {str(e)}")
    
    # If we get here, all methods failed
    logger.error("All PDF extraction methods failed")
    return "Could not extract text from the PDF. The file may be encrypted, scanned, or damaged."

# File extraction functions - add after the existing PDF functions
def extract_text_from_excel(file_path):
    """Extract text from Excel files"""
    try:
        if pandas_available:
            # Use pandas to read Excel file
            df_list = pd.read_excel(file_path, sheet_name=None)
            text = ""
            for sheet_name, df in df_list.items():
                text += f"\n== Sheet: {sheet_name} ==\n"
                text += df.to_string(index=False) + "\n\n"
            return text
        else:
            return "Excel processing requires pandas library which is not installed."
    except Exception as e:
        logger.warning(f"Excel extraction error: {str(e)}")
        return f"Error extracting Excel content: {str(e)}"

def extract_text_from_csv(file_path):
    """Extract text from CSV files"""
    try:
        text = ""
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            csv_reader = csv.reader(file)
            for row in csv_reader:
                text += " | ".join(row) + "\n"
        return text
    except Exception as e:
        logger.warning(f"CSV extraction error: {str(e)}")
        return f"Error extracting CSV content: {str(e)}"

def extract_text_from_docx(file_path):
    """Extract text from Word documents"""
    try:
        if docx_available:
            doc = docx.Document(file_path)
            text = ""
            # Extract text from paragraphs
            for para in doc.paragraphs:
                text += para.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text)
                    text += " | ".join(row_text) + "\n"
            return text
        else:
            return "Word document processing requires python-docx library which is not installed."
    except Exception as e:
        logger.warning(f"DOCX extraction error: {str(e)}")
        return f"Error extracting Word document content: {str(e)}"

def extract_text_from_pptx(file_path):
    """Extract text from PowerPoint presentations"""
    try:
        if pptx_available:
            prs = Presentation(file_path)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"\n=== Slide {i+1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        else:
            return "PowerPoint processing requires python-pptx library which is not installed."
    except Exception as e:
        logger.warning(f"PPTX extraction error: {str(e)}")
        return f"Error extracting PowerPoint content: {str(e)}"

def extract_text_from_image(file_path):
    """Extract text from images using OCR"""
    try:
        if tesseract_available:
            # Check if tesseract is installed and available in path
            try:
                version = pytesseract.get_tesseract_version()
                logger.info(f"Using Tesseract OCR version: {version}")
                image = Image.open(file_path)
                text = pytesseract.image_to_string(image)
                return text or "No text detected in image."
            except pytesseract.TesseractNotFoundError:
                return "Tesseract OCR executable not found. Please ensure Tesseract is installed and in your PATH."
        else:
            return "Image OCR requires pytesseract which is not installed."
    except Exception as e:
        logger.warning(f"Image OCR error: {str(e)}")
        return f"Error extracting text from image: {str(e)}"

def extract_text_from_file(file_path, file_type):
    """Extract text from various file types"""
    # Determine file type if not provided
    if not file_type:
        file_type = os.path.splitext(file_path)[1].lower()
    
    logger.info(f"Attempting to extract text from file type: {file_type}")
    
    # Extract based on file type
    if file_type in ['.pdf', 'application/pdf']:
        return extract_text_from_pdf(file_path)
    elif file_type in ['.xlsx', '.xls', 'application/vnd.ms-excel', 
                      'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet']:
        return extract_text_from_excel(file_path)
    elif file_type in ['.csv', 'text/csv']:
        return extract_text_from_csv(file_path)
    elif file_type in ['.docx', '.doc', 'application/msword',
                      'application/vnd.openxmlformats-officedocument.wordprocessingml.document']:
        return extract_text_from_docx(file_path)
    elif file_type in ['.pptx', '.ppt', 'application/vnd.ms-powerpoint',
                      'application/vnd.openxmlformats-officedocument.presentationml.presentation']:
        return extract_text_from_pptx(file_path)
    elif file_type in ['.png', '.jpg', '.jpeg', 'image/png', 'image/jpeg']:
        return extract_text_from_image(file_path)
    else:
        # Try to read as plain text
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                return file.read()
        except:
            return f"Unsupported file type: {file_type}"

# API Routes
@app.route('/', defaults={'path': ''})
@app.route('/<path:path>')
def serve(path):
    if path != "" and os.path.exists(app.static_folder + '/' + path): # Fixed comparison
        return send_from_directory(app.static_folder, path)
    else:
        return send_from_directory(app.static_folder, 'index.html')

@app.route("/api/health")
def health_check():
    """Health check endpoint for the API"""
    if db is None:  # Fixed comparison
        return jsonify({
            'status': 'error',
            'message': 'Database connection is not available',
            'timestamp': datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S'),
        }), 500
        
    try:
        # Try to ping the database
        db.command('ping')
        return jsonify({
            'status': 'healthy',
            'database': 'connected',
            'timestamp': datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S'),
            'environment': {
                'variables': {
                    var: bool(os.getenv(var)) for var in [
                        'GROQ_API_KEY',
                        'EMAIL_SENDER',
                        'EMAIL_PASSWORD',
                        'MONGO_URL',
                        'MONGO_URI'
                    ]
                }
            }
        })
    except Exception as e:
        logger.error(f"MongoDB health check failed: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Database connection error',
            'error': str(e),
            'timestamp': datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S')
        }), 500

@app.route("/api/chat", methods=["POST"])
def chat():
    """Send a message to the AI and get a response"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500

    data = request.get_json()
    user_message = data.get("message", "")
    model_name = data.get("model", "LLaMA3")
    language = data.get("language", "")
    chat_id = data.get("chatId", "")
    user_id = data.get("userId", "")

    # Auto-detect language if not specified
    if not language and user_message:
        try:
            detected_lang = detect(user_message)
            # Map detected language to our supported codes
            language = {
                'en': 'en',
                'hi': 'hi',
                'kn': 'kn',
                'ta': 'ta',
                'te': 'te',
                'sa': 'sa'
            }.get(detected_lang, 'en')
        except Exception as e:
            logger.warning(f"Language detection failed: {str(e)}")
            language = 'en'

    model = MODEL_MAPPING.get(model_name, MODEL_MAPPING["LLaMA3"])
    
    try:
        response = get_groq_response(user_message, language, model)
        
        # Save message to chat history if user is authenticated and valid chat ID provided
        if user_id and chat_id:
            try:
                # Add user message
                user_msg = {
                    "role": "user",
                    "content": user_message,
                    "language": language,
                    "timestamp": datetime.now()
                }
                
                # Add bot response
                bot_msg = {
                    "role": "bot",
                    "content": response,
                    "language": language,
                    "timestamp": datetime.now()
                }
                
                # Check if chat exists
                chat = db.chats.find_one({"_id": chat_id, "user_id": user_id})
                
                if chat:
                    # Update existing chat with new messages
                    db.chats.update_one(
                        {"_id": chat_id, "user_id": user_id},
                        {
                            "$push": {"messages": {"$each": [user_msg, bot_msg]}},
                            "$set": {"updated_at": datetime.now()}
                        }
                    )
                else:
                    # Create new chat with initial messages
                    # Auto-generate title from the first message
                    words = user_message.split()
                    auto_title = " ".join(words[:3]) + ("..." if len(words) > 3 else "")
                    title = auto_title.capitalize()
                    
                    db.chats.insert_one({
                        "_id": chat_id,
                        "user_id": user_id,
                        "title": title,
                        "messages": [user_msg, bot_msg],
                        "created_at": datetime.now(),
                        "updated_at": datetime.now()
                    })
                logger.info(f"Saved chat messages for user {user_id}, chat {chat_id}")
            except Exception as db_error:
                logger.error(f"Error saving chat messages: {str(db_error)}")
                # Continue to return response even if DB operation fails
        return jsonify({"reply": response})
    except Exception as e:
        logger.error(f"Chat error: {str(e)}")
        return jsonify({"reply": get_error_message(language)}), 500

@app.route("/api/signup", methods=["POST"])
def signup():
    """Direct signup without email verification"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    email = data.get("email", "").lower()
    password = data.get("password", "")
    full_name = data.get("fullName", "")
    phone_number = data.get("phoneNumber", "")

    if not email or not password:
        return jsonify({"success": False, "message": "Missing required fields"}), 400

    if '@' not in email:
        return jsonify({"success": False, "message": "Invalid email format"}), 400

    try:
        # Check if user already exists
        existing_user = db.users.find_one({"email": email})
        if existing_user:
            return jsonify({"success": False, "message": "Email already registered"}), 400
        
        # Use email username as name if not provided
        if not full_name:
            full_name = email.split('@')[0]
            
        # Create new user
        user = {
            "name": full_name,
            "email": email,
            "password": password,  # In production, this should be hashed
            "phone_number": phone_number,  # Store phone number
            "created_at": datetime.now()
        }
        
        result = db.users.insert_one(user)
        user_id = str(result.inserted_id)
        
        logger.info(f"User created with ID: {user_id}, email: {email}")
        
        return jsonify({
            "success": True,
            "user": {
                "id": user_id,
                "name": full_name,
                "email": email,
                "phoneNumber": phone_number
            }
        })
    except Exception as e:
        logger.error(f"Signup error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during signup: {str(e)}"}), 500

@app.route("/api/register", methods=["POST"])
def register():
    """Register with email verification"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    email = data.get("email", "").lower()
    full_name = data.get("fullName", "")
    password = data.get("password", "")
    phone_number = data.get("phoneNumber", "")

    logger.info(f"Registration attempt for {email} with name: {full_name}")

    if not email or not password:
        return jsonify({"success": False, "message": "Missing required fields"}), 400

    if '@' not in email:
        return jsonify({"success": False, "message": "Invalid email format"}), 400

    try:
        # Check if user already exists
        existing_user = db.users.find_one({"email": email})
        if existing_user:
            return jsonify({"success": False, "message": "Email already registered"}), 400
        
        # Use the email username as fullName if not provided
        if not full_name:
            full_name = email.split('@')[0]
        
        # Generate verification code
        code = generate_verification_code()
        expires_at = datetime.now() + timedelta(minutes=10)
        
        # Delete any existing OTPs for this email first
        db.verifications.delete_many({"email": email})
        
        # Store verification code in the database
        verification = {
            "email": email,
            "code": code,
            "full_name": full_name,
            "password": password,
            "phone_number": phone_number,  # Store phone number in verification
            "created_at": datetime.now(),
            "expires_at": expires_at
        }
        
        db.verifications.insert_one(verification)
        
        logger.info(f"Verification code for {email}: {code}")
        
        # Send verification email
        if send_verification_email(email, code):
            return jsonify({
                "success": True,
                "message": "Verification code sent to your email"
            })
        else:
            return jsonify({
                "success": False,
                "message": "Failed to send verification email"
            }), 500
    except Exception as e:
        logger.error(f"Registration error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during registration: {str(e)}"}), 500

@app.route("/api/verify", methods=["POST"])
def verify():
    """Verify email with the provided code"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500

    data = request.get_json()
    email = data.get("email", "").lower()
    code = data.get("code", "")

    logger.info(f"Verification attempt for {email} with code {code}")

    if not email or not code:
        return jsonify({"success": False, "message": "Missing required fields"}), 400

    try:
        # Find the verification code
        verification = db.verifications.find_one({
            "email": email,
            "code": code,
            "expires_at": {"$gt": datetime.now()}  # Check if not expired
        })
        
        if not verification:
            expired_verification = db.verifications.find_one({
                "email": email,
                "code": code
            })
            
            if expired_verification:
                return jsonify({"success": False, "message": "Verification code has expired"}), 400
            else:
                return jsonify({"success": False, "message": "Invalid verification code"}), 400
        
        # Check if user already exists
        existing_user = db.users.find_one({"email": email})
        if existing_user:
            user_id = str(existing_user["_id"])
            name = existing_user["name"]
            phone_number = existing_user.get("phone_number", "")
        else:
            # Create new user
            user = {
                "name": verification["full_name"],
                "email": email,
                "password": verification["password"],  # In production, this should be hashed
                "phone_number": verification.get("phone_number", ""),  # Get phone number from verification
                "created_at": datetime.now()
            }
            
            result = db.users.insert_one(user)
            user_id = str(result.inserted_id)
            name = verification["full_name"]
            phone_number = verification.get("phone_number", "")  # Get phone number from verification
        
        # Delete the verification
        db.verifications.delete_one({"_id": verification["_id"]})
        
        logger.info(f"Email verified successfully for {email}")
        
        return jsonify({
            "success": True,
            "message": "Email verified successfully",
            "user": {
                "id": user_id,
                "name": name,
                "email": email,
                "phoneNumber": phone_number  # Use properly defined variable
            }
        })
    except Exception as e:
        logger.error(f"Verification error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during verification: {str(e)}"}), 500

@app.route("/api/login", methods=["POST"])
def login():
    """Login with email and password"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    email = data.get("email", "").lower()
    password = data.get("password", "")

    if not email or not password:
        return jsonify({"success": False, "message": "Missing required fields"}), 400

    try:
        # Find user by email
        user = db.users.find_one({"email": email})
        
        if not user:
            return jsonify({"success": False, "message": "Email not registered"}), 401
        
        # Check password (in production, you should use a proper password hash comparison)
        if user["password"] != password:
            return jsonify({"success": False, "message": "Incorrect password"}), 401

        logger.info(f"User logged in: {email}")
        
        return jsonify({
            "success": True,
            "message": "Login successful",
            "user": {
                "id": str(user["_id"]),
                "name": user.get("name", email.split('@')[0]),
                "email": user["email"]
            }
        })
    except Exception as e:
        logger.error(f"Login error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during login: {str(e)}"}), 500

@app.route("/api/send-verification", methods=["POST"])
def send_verification():
    """Send verification code to email"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500

    data = request.get_json()
    email = data.get("email", "").lower()

    if not email:
        return jsonify({"success": False, "message": "Email is required"}), 400

    try:
        # Check if user exists
        user = db.users.find_one({"email": email})
        
        if not user:
            return jsonify({"success": False, "message": "Email not registered"}), 404
        
        # Generate verification code
        code = generate_verification_code()
        expires_at = datetime.now() + timedelta(minutes=10)
        
        # Delete any existing verifications for this email
        db.verifications.delete_many({"email": email})
        
        # Store verification code
        verification = {
            "email": email,
            "code": code,
            "created_at": datetime.now(),
            "expires_at": expires_at
        }
        
        db.verifications.insert_one(verification)
        
        logger.info(f"Verification code sent to {email}: {code}")
        
        # Send verification email
        if send_verification_email(email, code):
            return jsonify({
                "success": True,
                "message": "Verification code sent to your email"
            })
        else:
            return jsonify({
                "success": False,
                "message": "Failed to send verification email"
            }), 500
    except Exception as e:
        logger.error(f"Send verification error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred while sending verification: {str(e)}"}), 500

@app.route("/api/chats", methods=["GET"])
def get_chats():
    """Get all chats for a user"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500

    user_id = request.args.get("userId", "")

    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400

    try:
        # Find all chats for the user
        chats_cursor = db.chats.find({"user_id": user_id}).sort("updated_at", -1)

        # Convert to the format expected by the frontend
        formatted_chats = {}
        for chat in chats_cursor:
            formatted_messages = []
            for msg in chat.get("messages", []):
                formatted_messages.append({
                    "role": msg["role"],
                    "content": msg["content"],
                    "language": msg.get("language", "en"),
                    "timestamp": msg["timestamp"].isoformat() if isinstance(msg["timestamp"], datetime) else msg["timestamp"]
                })
            
            formatted_chats[chat["_id"]] = {
                "title": chat["title"],
                "messages": formatted_messages
            }
        logger.info(f"Retrieved {len(formatted_chats)} chats for user {user_id}")
        
        return jsonify({
            "success": True,
            "chats": formatted_chats
        })
    except Exception as e:
        logger.error(f"Get chats error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred while retrieving chats: {str(e)}"}), 500

@app.route("/api/chats", methods=["POST"])
def create_chat():
    """Create a new chat"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    user_id = data.get("userId", "")
    chat_id = data.get("chatId", f"chat_{datetime.now().timestamp()}")
    title = data.get("title", "New Chat")

    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400

    try:
        # Check if chat ID already exists
        existing_chat = db.chats.find_one({"_id": chat_id})
        if existing_chat:
            # Generate a new unique chat ID
            chat_id = f"chat_{datetime.now().timestamp()}_{random.randint(1000, 9999)}"

        # Create new chat
        chat = {
            "_id": chat_id,
            "user_id": user_id,
            "title": title,
            "messages": [],
            "created_at": datetime.now(),
            "updated_at": datetime.now()
        }
        
        db.chats.insert_one(chat)
        logger.info(f"Created new chat {chat_id} for user {user_id}")
        
        return jsonify({
            "success": True,
            "chatId": chat_id,
            "chat": {
                "title": title,
                "messages": []
            }
        })
    except Exception as e:
        logger.error(f"Create chat error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred while creating chat: {str(e)}"}), 500

@app.route("/api/chats/<chat_id>", methods=["PUT"])
def update_chat(chat_id):
    """Update a chat's title"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    user_id = data.get("userId", "")
    title = data.get("title", "")

    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400

    try:
        # Find the chat
        chat = db.chats.find_one({"_id": chat_id, "user_id": user_id})
        if not chat:
            return jsonify({"success": False, "message": "Chat not found"}), 404

        # Update the title if provided
        if title:
            db.chats.update_one(
                {"_id": chat_id, "user_id": user_id},
                {"$set": {"title": title, "updated_at": datetime.now()}}
            )
            logger.info(f"Updated title of chat {chat_id} to '{title}'")
        
        # Get the updated chat
        updated_chat = db.chats.find_one({"_id": chat_id, "user_id": user_id})
        
        # Format messages for response
        formatted_messages = []
        for msg in updated_chat.get("messages", []):
            formatted_messages.append({
                "role": msg["role"],
                "content": msg["content"],
                "language": msg.get("language", "en"),
                "timestamp": msg["timestamp"].isoformat() if isinstance(msg["timestamp"], datetime) else msg["timestamp"]
            })
        
        return jsonify({
            "success": True,
            "chat": {
                "title": updated_chat["title"],
                "messages": formatted_messages
            }
        })
    except Exception as e:
        logger.error(f"Update chat error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred while updating chat: {str(e)}"}), 500

@app.route("/api/chats/<chat_id>", methods=["DELETE"])
def delete_chat(chat_id):
    """Delete a chat"""
    if db is None:  # Fixed comparison
        return jsonify({"error": "Database connection is not available"}), 500

    user_id = request.args.get("userId", "")

    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400

    try:
        # Find and delete the chat
        result = db.chats.delete_one({"_id": chat_id, "user_id": user_id})
        if result.deleted_count == 0:
            return jsonify({"success": False, "message": "Chat not found"}), 404
        logger.info(f"Deleted chat {chat_id} for user {user_id}")
        
        return jsonify({
            "success": True,
            "message": "Chat deleted successfully"
        })
    except Exception as e:
        logger.error(f"Delete chat error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred while deleting chat: {str(e)}"}), 500

@app.route("/api/speech-to-text", methods=["POST"])
def speech_to_text():
    """Convert speech to text"""
    if not speech_recognition_available:
        return jsonify({"error": "Speech recognition is not available on the server"}), 503

    if 'audio' not in request.files:
        return jsonify({"error": "No audio file provided"}), 400

    audio_file = request.files['audio']
    if not audio_file.filename:
        return jsonify({"error": "Empty audio file"}), 400

    logger.info(f"Received audio file: {audio_file.filename}, size: {audio_file.content_length} bytes")

    try:
        # Create temp directory if it doesn't exist with proper permissions
        create_directory_with_permissions("temp")
        
        # Save the audio file temporarily with a unique name
        temp_filename = f"temp_audio_{datetime.now().timestamp()}_{random.randint(1000, 9999)}.wav"
        temp_path = os.path.join("temp", temp_filename)
        audio_file.save(temp_path)
        logger.info(f"Saved audio file to: {temp_path}")

        preferred_lang = request.form.get('language', '').lower()
        
        with sr.AudioFile(temp_path) as source:
            logger.info("Attempting to recognize speech...")
            audio_data = recognizer.record(source)

            # If preferred language is specified, try that first
            if preferred_lang and preferred_lang in SPEECH_LANGUAGES:
                try:
                    text = recognizer.recognize_google(audio_data, language=preferred_lang)
                    if text:
                        logger.info(f"Detected language: {preferred_lang}, text: {text}")
                        return jsonify({
                            "text": text,
                            "language": preferred_lang
                        })
                except Exception as e:
                    logger.warning(f"Recognition failed for preferred language {preferred_lang}: {str(e)}")

            # Fallback to trying all supported languages
            detected_text = None
            detected_lang = "en"  # default fallback

            # Try Kannada first if not already tried
            if preferred_lang != "kn" and "kn" in SPEECH_LANGUAGES:
                try:
                    text = recognizer.recognize_google(audio_data, language="kn")
                    if text:
                        detected_text = text
                        detected_lang = "kn"
                        logger.info(f"Detected language: kn, text: {text}")
                except Exception as e:
                    logger.warning(f"Recognition failed for Kannada: {str(e)}")

            if not detected_text:
                for lang in SPEECH_LANGUAGES:
                    if lang == "kn":  # Already tried
                        continue
                    try:
                        text = recognizer.recognize_google(audio_data, language=lang)
                        if text:
                            detected_text = text
                            detected_lang = lang
                            logger.info(f"Detected language: {lang}, text: {text}")
                            break
                    except Exception as e:
                        logger.warning(f"Recognition failed for language {lang}: {str(e)}")

            if detected_text:
                return jsonify({
                    "text": detected_text,
                    "language": detected_lang
                })

            logger.warning("Could not recognize speech in any supported language")
            return jsonify({"error": "Could not recognize speech"}), 400
    except Exception as e:
        logger.error(f"Speech recognition error: {str(e)}")
        return jsonify({"error": str(e)}), 500
    finally:
        # Clean up temp file
        try:
            if os.path.exists(temp_path):
                os.remove(temp_path)
                logger.info(f"Removed temporary file: {temp_path}")
        except Exception as cleanup_error:
            logger.warning(f"Failed to remove temporary file: {str(cleanup_error)}")

@app.route("/api/text-to-speech", methods=["POST"])
def text_to_speech_endpoint():
    """Convert text to speech"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({"error": "Invalid JSON data"}), 400

        text = data.get("text", "")
        language = data.get("language", "en")
    
        if not text:
            return jsonify({"error": "No text provided"}), 400

        # Limit text length for performance reasons
        if len(text) > 5000:
            text = text[:5000]
            logger.warning(f"Text truncated to 5000 characters")
    
        # Check if language is supported
        if language not in TTS_LANGUAGES:
            logger.warning(f"Language {language} not supported, falling back to English")
            language = "en"
    
        logger.info(f"Converting text to speech in {language}. Text length: {len(text)} characters")
    
        audio_path = text_to_speech(text, language)
        if audio_path:
            logger.info(f"Successfully generated speech, sending file: {audio_path}")
            return send_from_directory(os.path.dirname(audio_path), os.path.basename(audio_path))
        else:
            logger.error("Failed to generate speech audio")
            return jsonify({"error": "Could not generate speech"}), 500
    except Exception as e:
        logger.error(f"Text-to-speech error: {str(e)}")
        return jsonify({"error": f"An error occurred during speech synthesis: {str(e)}"}), 500

@app.route("/api/summarize-pdf", methods=["POST"])
def summarize_pdf():
    """Endpoint to summarize uploaded file content"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500

    if 'pdf' not in request.files:
        return jsonify({"success": False, "message": "No file provided"}), 400

    uploaded_file = request.files['pdf']

    if uploaded_file.filename == '':
        return jsonify({"success": False, "message": "Empty file"}), 400

    # Check file size (limit to 10MB)
    if uploaded_file.content_length and uploaded_file.content_length > 10 * 1024 * 1024:
        return jsonify({"success": False, "message": "File size exceeds 10MB limit"}), 400

    try:
        # Create temp directory if it doesn't exist with proper permissions
        create_directory_with_permissions("temp")

        # Get file extension
        file_ext = os.path.splitext(uploaded_file.filename)[1].lower()
        content_type = uploaded_file.content_type
        
        # Save the file to a temporary location
        temp_file_name = f"temp_file_{datetime.now().timestamp()}_{random.randint(1000, 9999)}{file_ext}"
        temp_file_path = os.path.join("temp", temp_file_name)
        uploaded_file.save(temp_file_path)
        logger.info(f"Saved temporary file to {temp_file_path}, type: {content_type}")

        # Extract text based on file type
        try:
            full_text = extract_text_from_file(temp_file_path, content_type)
            if not full_text or not full_text.strip():
                return jsonify({"success": False, "message": "Could not extract any text from the file. The file might be empty, encrypted, or contains no readable text."}), 400
            logger.info(f"Extracted {len(full_text)} characters from file")
        except Exception as extraction_error:
            logger.error(f"Error extracting text from file: {str(extraction_error)}")
            return jsonify({"success": False, "message": f"Failed to extract text from file: {str(extraction_error)}"}), 500
        finally:
            # Clean up temporary file
            try:
                if os.path.exists(temp_file_path):
                    os.remove(temp_file_path)
                    logger.info(f"Removed temporary file: {temp_file_path}")
            except Exception as cleanup_error:
                logger.warning(f"Failed to remove temporary file: {str(cleanup_error)}")

        # Initialize content collection if it doesn't exist
        if "file_contents" not in db.list_collection_names():
            db.create_collection("file_contents")
            logger.info("Created file_contents collection")

        # Store extracted text in MongoDB collection
        # Limit text length to avoid MongoDB document size limits
        max_text_length = 500000  # MongoDB documents have a 16MB limit
        truncated_text = full_text[:max_text_length] if len(full_text) > max_text_length else full_text
        file_content_doc = {
            "filename": uploaded_file.filename,
            "content_type": content_type,
            "content": truncated_text,
            "uploaded_at": datetime.utcnow()
        }
        result = db.file_contents.insert_one(file_content_doc)
        file_content_id = str(result.inserted_id)

        # Use Groq API to summarize the extracted text
        # Limit prompt size to avoid token limits
        max_prompt_size = 4000
        text_to_summarize = truncated_text[:max_prompt_size] if len(truncated_text) > max_prompt_size else truncated_text
        
        try:
            prompt = f"Summarize the following text extract from a file. Provide a clear, concise overview that tells what the file is about, including key topics, structure, and important insights:\n\n{text_to_summarize}"
            summary = groq_client.chat.completions.create(
                model="llama3-8b-8192",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.5,
                max_tokens=512,
                top_p=1,
                frequency_penalty=0,
                presence_penalty=0,
                stop=None,
            ).choices[0].message.content
            logger.info(f"Generated summary of length {len(summary)}")
        except Exception as api_error:
            logger.error(f"Error calling Groq API: {str(api_error)}")
            summary = "Failed to generate summary. You can still ask questions about the document below."

        return jsonify({
            "success": True,
            "summary": summary,
            "fileContentId": file_content_id,
            "fileName": uploaded_file.filename,
            "fileType": content_type
        })
    except Exception as e:
        logger.error(f"Error summarizing file: {str(e)}")
        return jsonify({
            "success": False,
            "message": f"Failed to process file: {str(e)}"
        }), 500

@app.route("/api/ask-pdf-question", methods=["POST"])
def ask_pdf_question():
    """Endpoint to ask questions based on uploaded file content"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500

    data = request.get_json()
    file_content_id = data.get("pdfContentId", "")  # Keep param name for backward compatibility
    question = data.get("question", "")

    if not file_content_id or not question:
        return jsonify({"success": False, "message": "Missing fileContentId or question"}), 400

    try:
        # Retrieve file content from MongoDB
        try:
            # Look in both collections for backward compatibility
            file_doc = db.file_contents.find_one({"_id": ObjectId(file_content_id)})
            if not file_doc:
                file_doc = db.pdf_contents.find_one({"_id": ObjectId(file_content_id)})
                
        except Exception as db_error:
            logger.error(f"Error retrieving file content: {str(db_error)}")
            return jsonify({"success": False, "message": "Invalid file content ID"}), 400

        if not file_doc:
            return jsonify({"success": False, "message": "File content not found"}), 404

        file_text = file_doc.get("content", "")
        file_name = file_doc.get("filename", "Document")
        content_type = file_doc.get("content_type", "text/plain")

        if not file_text.strip():
            return jsonify({"success": False, "message": "Stored file content is empty"}), 400

        # Use a larger model for better comprehension, if available
        model = "llama3-70b-8192" if "llama3-70b-8192" in MODEL_MAPPING.values() else "llama3-8b-8192"
        
        # Handle very short questions or greetings differently
        if len(question.strip()) < 5 or question.lower() in ["hi", "hello", "hey"]:
            prompt = f"""You are a helpful document assistant analyzing a {content_type} file named '{file_name}'.
            The user has just sent a greeting or very short message: "{question}"
            Please respond professionally, let them know you're ready to answer questions about the document,
            and provide a brief 1-2 sentence overview of what the document appears to be about."""
        else:
            # Limit context size to avoid token limits but keep enough context
            max_context_size = 6000  # Increased from 4000
            truncated_text = file_text[:max_context_size] if len(file_text) > max_context_size else file_text
            
            prompt = f"""You are a document analysis AI assistant. Your task is to answer questions based EXCLUSIVELY on the 
            information provided in the following {content_type} file named '{file_name}'.

            DOCUMENT CONTENT:
            ```
            {truncated_text}
            ```

            USER QUESTION: {question}

            Important instructions:
            1. Answer ONLY based on information in the document above.
            2. If the answer is not in the document, clearly state that you cannot find this information in the document.
            3. Do NOT make up information or use your general knowledge unless it directly relates to understanding the document.
            4. Provide specific references to where in the document you found the information when possible.
            5. Be concise but thorough in your answer.
            """

        logger.info(f"Sending document question to model {model} with prompt length {len(prompt)}")
        
        try:
            answer = groq_client.chat.completions.create(
                model=model,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,  # Lower temperature for more factual responses
                max_tokens=800,   # Increased to allow for more detailed answers
                top_p=1,
                frequency_penalty=0,
                presence_penalty=0,
                stop=None,
            ).choices[0].message.content
        except Exception as api_error:
            logger.error(f"Error calling Groq API: {str(api_error)}")
            answer = "I'm sorry, I encountered an error processing your question about this document. Please try again later."

        return jsonify({"success": True, "answer": answer})
    except Exception as e:
        logger.error(f"Error answering file question: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to answer question: {str(e)}"}), 500

@app.route("/api/reset-password", methods=["POST"])
def reset_password():
    """Reset user password"""
    if db is None:
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    user_id = data.get("userId")
    current_password = data.get("currentPassword")
    new_password = data.get("newPassword")

    if not user_id or not current_password or not new_password:
        return jsonify({"success": False, "message": "Missing required fields"}), 400

    if len(new_password) < 6:
        return jsonify({"success": False, "message": "New password must be at least 6 characters long"}), 400

    try:
        # Find user by ID
        try:
            user = db.users.find_one({"_id": ObjectId(user_id)})
        except Exception as e:
            logger.error(f"Error finding user: {str(e)}")
            return jsonify({"success": False, "message": "Invalid user ID"}), 400
        
        if not user:
            return jsonify({"success": False, "message": "User not found"}), 404
        
        # Check current password (in production, you should use a proper password hash comparison)
        if user["password"] != current_password:
            return jsonify({"success": False, "message": "Current password is incorrect"}), 401
        
        # Update password
        db.users.update_one(
            {"_id": ObjectId(user_id)},
            {"$set": {"password": new_password}}  # In production, this should be hashed
        )
        
        logger.info(f"Password reset successful for user {user_id}")
        
        return jsonify({
            "success": True,
            "message": "Password reset successful"
        })
    except Exception as e:
        logger.error(f"Password reset error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during password reset: {str(e)}"}), 500

@app.route("/api/update-profile", methods=["PUT"])
def update_profile():
    """Update user profile information"""
    if db is None:
        return jsonify({"error": "Database connection is not available"}), 500
        
    data = request.get_json()
    user_id = data.get("userId")
    name = data.get("name")

    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400

    try:
        # Find user by ID
        try:
            user = db.users.find_one({"_id": ObjectId(user_id)})
        except Exception as e:
            logger.error(f"Error finding user: {str(e)}")
            return jsonify({"success": False, "message": "Invalid user ID"}), 400
        
        if not user:
            return jsonify({"success": False, "message": "User not found"}), 404
        
        # Update fields that are provided
        update_fields = {}
        if name:
            update_fields["name"] = name
        
        if update_fields:
            db.users.update_one(
                {"_id": ObjectId(user_id)},
                {"$set": update_fields}
            )
        
        logger.info(f"Profile updated for user {user_id}")
        
        # Get updated user data
        updated_user = db.users.find_one({"_id": ObjectId(user_id)})
        
        return jsonify({
            "success": True,
            "message": "Profile updated successfully",
            "user": {
                "id": user_id,
                "name": updated_user.get("name"),
                "email": updated_user.get("email"),
                "phoneNumber": updated_user.get("phone_number", "")  # Use updated_user instead of user
            }
        })
    except Exception as e:
        logger.error(f"Profile update error: {str(e)}")
        return jsonify({"success": False, "message": f"An error occurred during profile update: {str(e)}"}), 500

@app.errorhandler(404)
def not_found(error):
    return jsonify({"error": "Endpoint not found"}), 404

@app.errorhandler(405)
def method_not_allowed(error):
    return jsonify({"error": "Method not allowed"}), 405

@app.errorhandler(500)
def internal_server_error(error):
    logger.error(f"Internal server error: {str(error)}")
    return jsonify({"error": f"Internal server error: {str(error)}"}), 500

# Import new modules for model training
import training
import model_utils

# New routes for model training system

@app.route("/api/training/upload", methods=["POST"])
def upload_training_data():
    """Upload training data for model training"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    if 'file' not in request.files:
        return jsonify({"success": False, "message": "No file provided"}), 400
    
    user_id = request.form.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    data_name = request.form.get("name", "Untitled Dataset")
    data_description = request.form.get("description", "")
    data_format = request.form.get("format", "raw")  # raw, classification, entity, etc.
    
    uploaded_file = request.files['file']
    if uploaded_file.filename == '':
        return jsonify({"success": False, "message": "Empty file"}), 400
    
    # Check file size (limit to 50MB)
    if uploaded_file.content_length and uploaded_file.content_length > 50 * 1024 * 1024:
        return jsonify({"success": False, "message": "File size exceeds 50MB limit"}), 400
    
    try:
        # Create training_data directory if it doesn't exist
        create_directory_with_permissions("training_data")
        
        # Save the file with a unique name
        file_ext = os.path.splitext(uploaded_file.filename)[1].lower()
        content_type = uploaded_file.content_type
        file_name = f"training_{user_id}_{datetime.now().timestamp()}{file_ext}"
        file_path = os.path.join("training_data", file_name)
        
        uploaded_file.save(file_path)
        logger.info(f"Saved training data file to {file_path}")
        
        # Create entry in training_data collection
        training_data = {
            "name": data_name,
            "description": data_description,
            "user_id": user_id,
            "file_path": file_path,
            "file_type": file_ext if file_ext else content_type,
            "data_format": data_format,
            "created_at": datetime.now(),
            "processed": False
        }
        
        result = db.training_data.insert_one(training_data)
        data_id = str(result.inserted_id)
        
        # Process the training data in the background
        # In production, you would use a task queue like Celery
        # For simplicity, we'll do it synchronously here
        processing_success = training.process_training_data(data_id)
        
        return jsonify({
            "success": True,
            "message": "Training data uploaded successfully",
            "dataId": data_id,
            "processed": processing_success
        })
    except Exception as e:
        logger.error(f"Error uploading training data: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to upload training data: {str(e)}"}), 500

@app.route("/api/training/data", methods=["GET"])
def get_training_data():
    """Get all training data for a user"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.args.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find all training data for the user
        data_cursor = db.training_data.find({"user_id": user_id}).sort("created_at", -1)
        
        training_data = []
        for data in data_cursor:
            training_data.append({
                "id": str(data["_id"]),
                "name": data["name"],
                "description": data["description"],
                "fileType": data["file_type"],
                "format": data["data_format"],
                "createdAt": data["created_at"].isoformat(),
                "processed": data["processed"],
                "rowCount": data.get("row_count", "Unknown")
            })
        
        return jsonify({
            "success": True,
            "trainingData": training_data
        })
    except Exception as e:
        logger.error(f"Error getting training data: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to get training data: {str(e)}"}), 500

@app.route("/api/training/config", methods=["POST"])
def create_model_config():
    """Create a new model configuration"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "No data provided"}), 400
    
    user_id = data.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    name = data.get("name", "Untitled Model")
    description = data.get("description", "")
    model_type = data.get("modelType")
    base_model = data.get("baseModel")
    hyperparameters = data.get("hyperparameters", {})
    
    if not model_type or not base_model:
        return jsonify({"success": False, "message": "Model type and base model are required"}), 400
    
    try:
        # Create model config
        config = {
            "name": name,
            "description": description,
            "user_id": user_id,
            "model_type": model_type,
            "base_model": base_model,
            "hyperparameters": json.dumps(hyperparameters),
            "created_at": datetime.now(),
            "updated_at": datetime.now()
        }
        
        result = db.model_configs.insert_one(config)
        config_id = str(result.inserted_id)
        
        return jsonify({
            "success": True,
            "message": "Model configuration created successfully",
            "configId": config_id
        })
    except Exception as e:
        logger.error(f"Error creating model config: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to create model config: {str(e)}"}), 500

@app.route("/api/training/config", methods=["GET"])
def get_model_configs():
    """Get all model configurations for a user"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.args.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find all model configs for the user
        config_cursor = db.model_configs.find({"user_id": user_id}).sort("created_at", -1)
        
        configs = []
        for config in config_cursor:
            configs.append({
                "id": str(config["_id"]),
                "name": config["name"],
                "description": config["description"],
                "modelType": config["model_type"],
                "baseModel": config["base_model"],
                "hyperparameters": json.loads(config["hyperparameters"]) if config["hyperparameters"] else {},
                "createdAt": config["created_at"].isoformat()
            })
        
        return jsonify({
            "success": True,
            "modelConfigs": configs
        })
    except Exception as e:
        logger.error(f"Error getting model configs: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to get model configs: {str(e)}"}), 500

@app.route("/api/training/job", methods=["POST"])
def create_training_job():
    """Create a new model training job"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "No data provided"}), 400
    
    user_id = data.get("userId")
    config_id = data.get("configId")
    data_ids = data.get("dataIds", [])
    name = data.get("name", "Training Job")
    
    if not user_id or not config_id or not data_ids:
        return jsonify({
            "success": False, 
            "message": "User ID, config ID, and at least one data ID are required"
        }), 400
    
    try:
        # Create training job
        job_id = training.create_training_job(config_id, data_ids, user_id, name)
        if not job_id:
            return jsonify({
                "success": False,
                "message": "Failed to create training job. Check logs for details."
            }), 500
        
        # Start training in the background
        # In production, you would use a task queue like Celery
        # For simplicity, we'll do it in a separate thread
        import threading
        training_thread = threading.Thread(target=training.run_training_job, args=(job_id,))
        training_thread.daemon = True
        training_thread.start()
        
        return jsonify({
            "success": True,
            "message": "Training job created and started successfully",
            "jobId": job_id
        })
    except Exception as e:
        logger.error(f"Error creating training job: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to create training job: {str(e)}"}), 500

@app.route("/api/training/job", methods=["GET"])
def get_training_jobs():
    """Get all training jobs for a user"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.args.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find all training jobs for the user
        job_cursor = db.model_training_jobs.find({"user_id": user_id}).sort("created_at", -1)
        
        jobs = []
        for job in job_cursor:
            jobs.append({
                "id": str(job["_id"]),
                "name": job["name"],
                "configId": job["config_id"],
                "status": job["status"],
                "createdAt": job["created_at"].isoformat(),
                "startedAt": job["started_at"].isoformat() if job.get("started_at") else None,
                "completedAt": job["completed_at"].isoformat() if job.get("completed_at") else None,
                "resultModelId": job.get("result_model_id"),
                "logs": job.get("logs", [])
            })
        
        return jsonify({
            "success": True,
            "trainingJobs": jobs
        })
    except Exception as e:
        logger.error(f"Error getting training jobs: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to get training jobs: {str(e)}"}), 500

@app.route("/api/training/job/<job_id>", methods=["GET"])
def get_training_job(job_id):
    """Get details of a specific training job"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.args.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find the training job
        job = db.model_training_jobs.find_one({
            "_id": ObjectId(job_id),
            "user_id": user_id
        })
        
        if not job:
            return jsonify({"success": False, "message": "Training job not found"}), 404
        
        # Get the config details
        config = db.model_configs.find_one({"_id": ObjectId(job["config_id"])})
        config_details = {
            "id": str(config["_id"]),
            "name": config["name"],
            "modelType": config["model_type"],
            "baseModel": config["base_model"]
        } if config else None
        
        # Get training data details
        training_data_details = []
        for data_id in job.get("training_data_ids", []):
            data = db.training_data.find_one({"_id": ObjectId(data_id)})
            if data:
                training_data_details.append({
                    "id": str(data["_id"]),
                    "name": data["name"],
                    "format": data["data_format"]
                })
        
        # Get result model if available
        model_details = None
        if job.get("result_model_id"):
            model = db.trained_models.find_one({"_id": ObjectId(job["result_model_id"])})
            if model:
                try:
                    metrics = json.loads(model["metrics"]) if model["metrics"] else {}
                except:
                    metrics = {}
                    
                model_details = {
                    "id": str(model["_id"]),
                    "name": model["name"],
                    "version": model["version"],
                    "metrics": metrics,
                    "isActive": model["is_active"]
                }
        
        return jsonify({
            "success": True,
            "job": {
                "id": str(job["_id"]),
                "name": job["name"],
                "status": job["status"],
                "createdAt": job["created_at"].isoformat(),
                "startedAt": job["started_at"].isoformat() if job.get("started_at") else None,
                "completedAt": job["completed_at"].isoformat() if job.get("completed_at") else None,
                "logs": job.get("logs", []),
                "config": config_details,
                "trainingData": training_data_details,
                "resultModel": model_details
            }
        })
    except Exception as e:
        logger.error(f"Error getting training job: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to get training job: {str(e)}"}), 500

@app.route("/api/models/list", methods=["GET"])
def get_user_models():
    """Get all trained models for a user"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.args.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find all models for the user
        model_cursor = db.trained_models.find({"user_id": user_id}).sort("created_at", -1)
        
        models = []
        for model in model_cursor:
            try:
                metrics = json.loads(model["metrics"]) if model["metrics"] else {}
            except:
                metrics = {}
                
            models.append({
                "id": str(model["_id"]),
                "name": model["name"],
                "description": model["description"],
                "modelType": model["model_type"],
                "baseModel": model["base_model"],
                "version": model["version"],
                "metrics": metrics,
                "createdAt": model["created_at"].isoformat(),
                "isActive": model["is_active"]
            })
        
        return jsonify({
            "success": True,
            "models": models
        })
    except Exception as e:
        logger.error(f"Error getting user models: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to get user models: {str(e)}"}), 500

@app.route("/api/models/<model_id>/activate", methods=["POST"])
def activate_model(model_id):
    """Activate a trained model for use"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    user_id = request.json.get("userId")
    if not user_id:
        return jsonify({"success": False, "message": "User ID is required"}), 400
    
    try:
        # Find the model
        model = db.trained_models.find_one({
            "_id": ObjectId(model_id),
            "user_id": user_id
        })
        
        if not model:
            return jsonify({"success": False, "message": "Model not found"}), 404
        
        # Deactivate all other models of the same type
        db.trained_models.update_many(
            {
                "user_id": user_id,
                "model_type": model["model_type"],
                "_id": {"$ne": ObjectId(model_id)}
            },
            {"$set": {"is_active": False}}
        )
        
        # Activate this model
        db.trained_models.update_one(
            {"_id": ObjectId(model_id)},
            {"$set": {"is_active": True}}
        )
        
        return jsonify({
            "success": True,
            "message": f"Model {model['name']} activated successfully"
        })
    except Exception as e:
        logger.error(f"Error activating model: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to activate model: {str(e)}"}), 500

@app.route("/api/document/index", methods=["POST"])
def index_document_endpoint():
    """Index a document for semantic search"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "No data provided"}), 400
    
    document_id = data.get("documentId")
    text = data.get("text")
    metadata = data.get("metadata", {})
    
    if not document_id or not text:
        return jsonify({"success": False, "message": "Document ID and text are required"}), 400
    
    try:
        # Index the document
        result = model_utils.index_document(document_id, text, metadata)
        
        if result:
            return jsonify({
                "success": True,
                "message": "Document indexed successfully"
            })
        else:
            return jsonify({
                "success": False,
                "message": "Failed to index document. Check logs for details."
            }), 500
    except Exception as e:
        logger.error(f"Error indexing document: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to index document: {str(e)}"}), 500

@app.route("/api/document/search", methods=["POST"])
def search_documents_endpoint():
    """Search indexed documents using semantic similarity"""
    if db is None:
        return jsonify({"success": False, "message": "Database connection is not available"}), 500
    
    data = request.get_json()
    if not data:
        return jsonify({"success": False, "message": "No data provided"}), 400
    
    query = data.get("query")
    limit = data.get("limit", 5)
    
    if not query:
        return jsonify({"success": False, "message": "Search query is required"}), 400
    
    try:
        # Search documents
        results = model_utils.search_documents(query, limit)
        
        return jsonify({
            "success": True,
            "results": results
        })
    except Exception as e:
        logger.error(f"Error searching documents: {str(e)}")
        return jsonify({"success": False, "message": f"Failed to search documents: {str(e)}"}), 500

if __name__ == "__main__":
    # Define the port for the server
    port = int(os.environ.get("PORT", 5000))
    
    if db is None:
        logger.critical("MongoDB connection failed. Application may not function correctly.")
        
    try:
        from waitress import serve
        logger.info(f"Starting server with Waitress at {datetime.now().isoformat()} on port {port}")
        serve(app, host="0.0.0.0", port=port)
    except ImportError:
        logger.warning("Waitress not found, using Flask development server instead")
        app.run(host="0.0.0.0", port=port, debug=False)
    except Exception as e:
        logger.critical(f"Failed to start server: {str(e)}")
        print(f"Fatal error: {str(e)}")